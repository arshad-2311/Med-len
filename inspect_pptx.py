from pptx import Presentation
import sys

def inspect_pptx(file_path):
    prs = Presentation(file_path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for j, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()
            # Only print shapes that have some text to identify them
            if text:
                print(f"  Shape {j}: {text[:100].replace(chr(10), ' ')}")

if __name__ == "__main__":
    inspect_pptx(sys.argv[1])
