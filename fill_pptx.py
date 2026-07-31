from pptx import Presentation
import sys

def fill_pptx(input_path, output_path):
    prs = Presentation(input_path)
    
    # Slide 4
    for shape in prs.slides[3].shapes:
        if not shape.has_text_frame: continue
        text = shape.text.strip()
        if text.startswith("Team Name"):
            shape.text_frame.text = "Team Name : MedLens"
        elif text.startswith("Team Members"):
            shape.text_frame.text = "Team Members : Arshad Ahmed, Soham Mahangare"
        elif text.startswith("Problem Statement"):
            shape.text_frame.text = "Problem Statement : Medicine information is inaccessible to elders and non-English speakers, leading to dangerous errors."

    # Slide 5
    for shape in prs.slides[4].shapes:
        if not shape.has_text_frame: continue
        if "What everyday challenge" in shape.text:
            shape.text_frame.text = "MedLens aims to solve the critical issue of medication illiteracy. Elders and people who don't know English often struggle to read complex medical jargon, small text on medicine strips, and handwritten prescriptions. This lack of understanding leads to incorrect dosages, missed side effects, and serious health risks."

    # Slide 6
    for shape in prs.slides[5].shapes:
        if not shape.has_text_frame: continue
        if "What is your AI-powered idea?" in shape.text:
            shape.text_frame.text = "MedLens is an AI vernacular health companion. It uses computer vision (OCR) to scan medicine packages and prescriptions, and Generative AI to translate complex medical data into simple, easy-to-understand descriptions. Crucially, it translates these insights into regional Indian languages and offers an Audio Playback (TTS) feature so elders can simply listen to their medicine instructions."

    # Slide 7
    for shape in prs.slides[6].shapes:
        if not shape.has_text_frame: continue
        if "Who benefits from the solution?" in shape.text:
            shape.text_frame.text = "The primary beneficiaries are elderly patients and non-English speaking individuals in rural and semi-urban areas. Their main challenge is the inability to read or understand the medicines they are prescribed. MedLens empowers them to be independent, ensuring they know exactly what a drug does, how to take it safely, and any warnings, all in their native language."

    # Slide 8
    for shape in prs.slides[7].shapes:
        if not shape.has_text_frame: continue
        if "How would someone use your solution?" in shape.text:
            shape.text_frame.text = "1. User points their phone camera at a medicine strip or prescription.\n2. MedLens instantly scans and identifies the medicine.\n3. The AI breaks down the uses and side effects into plain language.\n4. The user selects their preferred language (e.g., Hindi) and taps 'Play' to hear the instructions read aloud. It is highly intuitive, requiring only a simple point-and-scan interaction."

    # Slide 9
    for shape in prs.slides[8].shapes:
        if not shape.has_text_frame: continue
        if "Where does AI contribute value?" in shape.text:
            shape.text_frame.text = "AI is the core engine of MedLens. Computer Vision extracts text from blurry, small, or handwritten medical sources. Generative AI (like Gemini) normalizes complex medical jargon into simple summaries. Text-To-Speech (TTS) AI converts these regional language summaries into natural audio. AI automates the entire process from scan to audio explanation."

    # Slide 10
    for shape in prs.slides[9].shapes:
        if not shape.has_text_frame: continue
        if "How can diverse users benefit?" in shape.text:
            shape.text_frame.text = "By providing multi-lingual support (Hindi, Marathi, Telugu, Tamil, etc.) and an audio-first interface, MedLens completely removes language and literacy barriers. The UI is designed to be highly accessible with large text, simple flows, and high-contrast elements, ensuring elders with poor eyesight can easily use the app."

    # Slide 11
    for shape in prs.slides[10].shapes:
        if not shape.has_text_frame: continue
        if "What positive change could this create?" in shape.text:
            shape.text_frame.text = "MedLens will drastically reduce medication errors, prevent adverse drug interactions, and empower patients to take control of their health. It can seamlessly integrate into broader ecosystems to provide wellness tools for employees and their families, creating a broader impact on community health and safety."

    # Slide 12
    for shape in prs.slides[11].shapes:
        if not shape.has_text_frame: continue
        if "How could the solution evolve?" in shape.text:
            shape.text_frame.text = "Future features include identifying generic or 'Jan Aushadhi' alternatives to save users money, providing dosage reminders via WhatsApp or SMS, and building a safety companion that warns of drug-drug interactions based on a user's medical history."

    prs.save(output_path)
    print(f"Filled PPTX saved to {output_path}")

if __name__ == "__main__":
    fill_pptx(sys.argv[1], sys.argv[2])
